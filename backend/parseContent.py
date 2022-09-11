import collections 
import collections.abc
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import csv
import PyPDF2


'''
two functions are implemented: 
pptToCsv : extract words content in .ppt to .csv, also extract .jpg and .png from .ppt.
pdfToCsv : extract words in .pdf to .csv   
    
'''


class pptToCsv(object):
    '''
    filePath   : the file path where .ppt source is, the file name need to be included.
    csvFileName: the csv name, user don't need to add .csv at the end.
    '''
    def __init__(self, pptPath ,folderPath, csvName) -> None:
        self.csvName = csvName
        self.folderPath = folderPath
        self.csvpath = f'{folderPath}{csvName}' 
        self.pptx = Presentation(pptPath)
    
    
 
    #each page in output csv file, will be seperated by '[pageX]'
    def main(self):
        # extract images        
        picPage = []
        def iter_picture_shapes(prs):
            pgNum = 0
            for slide in prs.slides:
                pgNum += 1
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE :
                        picPage.append(pgNum)
                        yield shape
        #turn img into .jpg or .png                
        count = 0
        imgNameList = []
        for picture in iter_picture_shapes(self.pptx):
            image = picture.image
            # ---get image "file" contents---
            image_bytes = image.blob
            # ---make up a name for the file, e.g. 'image.jpg'---
            image_filename = 'image_on_page%d.%s' % (picPage[count], image.ext)
            imgNameList.append([picPage[count],image_filename])
            count+=1
            with open(f'{self.folderPath}{image_filename}', 'wb') as f:
                f.write(image_bytes)

        #extract word content
        f = open(self.csvpath, 'w')
        writer = csv.writer(f)
        count = 1
        for slide in self.pptx.slides:
            writer.writerow([f'[page {count}]'])
            for shape in slide.shapes:
                if shape.has_text_frame :
                    textFrame = shape.text_frame
                    #write word content in .csv
                    for paragraph in textFrame.paragraphs:
                        str = paragraph.text
                        writer.writerow([str])
                    #write img html link in .csv
                    for item in imgNameList:
                        if item[0] == count:
                            str = f'<img src=\'{item[1]}\' />'
                            writer.writerow([str])
            count += 1        

        #output a txt of image name list
        with open(f'{self.folderPath}imageNameList.txt', 'w') as f:
            for item in imgNameList:
                name = item[1]
                f.write(f'{name}\n')
#to do list
# output a txt of images name list


class pdfToCsv(object):
    """
    folderPath : the folder path where the .csv will be stored.
    pdfPath    : the file path where .pdf source is, the file name need to be included.
    csvName    : the csv name, user don't need to add .csv at the end.
    """
    def __init__(self, folderPath, pdfPath, csvName) -> None:
        self.csvpath = f'{folderPath}{csvName}' 
        self.pdfPath = pdfPath
    
    
    #each page in output csv file, will be seperated by '#####'
    def main(self):
        f = open(f'{self.csvpath}.csv', 'w')
        writer = csv.writer(f)
        with open(self.pdfPath, "rb") as pdf_file:
            read_pdf = PyPDF2.PdfFileReader(pdf_file)
            count = '#####'
            for page in read_pdf.pages:
                writer.writerow([count])
                page_content = page.extractText()
                writer.writerow([page_content])
        
        f.close()

        
        

if __name__ == '__main__':
    pptToCsv = pptToCsv('/Users/adamliu/Desktop/semester1/internet technilogy/Week2-Lecture3.pptx',"/Users/adamliu/python/WIT hackthon 2022/UniMemoi/backend/","test")
    pptToCsv.main()
#     pdfToCsv = pdfToCsv("/Users/adamliu/python/WIT hackthon 2022/UniMemoi/backend/","/Users/adamliu/Desktop/semester1/algorithm and Complexity/01.pdf", "test")
#     pdfToCsv.main()