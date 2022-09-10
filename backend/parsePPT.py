import collections 
import collections.abc
from pptx import Presentation
import csv
from csv import reader

class pptToCsv(object):
    def __init__(self, filePath, csvFileName) -> None:
        self.csvFileName = csvFileName
        self.filePath =  filePath
        self.pptx = Presentation(self.filePath)

    def main(self):
        f = open("/Users/adamliu/python/WIT hackthon 2022/UniMemoi/{}.csv".format(self.csvFileName), 'w')
        writer = csv.writer(f)
        strList = []
        for slide in self.pptx.slides:
            for shape in slide.shapes:
                if shape.has_text_frame :
                    textFrame = shape.text_frame
                    for paragraph in textFrame.paragraphs:
                        str = paragraph.text
                        writer.writerow([str])


    
if __name__ == '__main__':
    pptToCsv = pptToCsv("/Users/adamliu/Desktop/semester1/internet technilogy/Week2-Lecture2.pptx", "test")
    pptToCsv.main()